from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {
        "title": "Galaxy Salon System",
        "bullets": [
            "Comprehensive Billing, Operations, and Analytics Platform",
            "Salon + Academy management in one application",
        ],
        "notes": [
            "Introduce scope: salon + academy + integrations.",
            "Mention stack: Next.js, Express, MongoDB.",
        ],
    },
    {
        "title": "Problem and Goal",
        "bullets": [
            "Pain points: slow billing, manual records, disconnected systems",
            "Goal: one integrated platform for counter, operations, and management",
            "Outcome: faster operations and cleaner business visibility",
        ],
        "notes": [
            "Focus on business impact first.",
            "Position this as an operations transformation project.",
        ],
    },
    {
        "title": "Who Can Use It",
        "bullets": [
            "Salon owners and managers",
            "Academy directors and instructors",
            "Front desk and billing staff",
            "Inventory and procurement personnel",
        ],
        "notes": [
            "Highlight the role-based access control that ensures each user sees only relevant modules.",
        ],
    },
    {
        "title": "Advantages Over Manual Systems",
        "bullets": [
            "Eliminates paper-based records and manual calculations",
            "Real-time inventory and financial updates",
            "Integrated appointment reminders via WhatsApp",
            "Data-driven decisions with reporting and AI insights",
        ],
        "notes": [
            "Emphasize time savings and error reduction.",
        ],
    },
    {
        "title": "Differentiators vs. Generic Billing Software",
        "bullets": [
            "Combines salon POS with academy management in one platform",
            "Built-in AI analytics for revenue prediction and customer insights",
            "Seamless integrations: Razorpay, WhatsApp, thermal printing, barcode scanning",
            "Role-based access control for staff and management",
        ],
        "notes": [
            "Most billing software focuses only on transactions; this system covers end-to-end operations.",
        ],
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Frontend: Next.js Pages Router + React",
            "Backend: Express API with JWT and validators",
            "Database: MongoDB with Mongoose",
            "Integrations: Razorpay, WhatsApp Cloud API, thermal print, barcode",
        ],
        "notes": [
            "Explain end-to-end request lifecycle from UI to DB.",
        ],
    },
    {
        "title": "Core Modules",
        "bullets": [
            "POS and Billing",
            "Appointments",
            "Customers and Loyalty",
            "Products and Inventory",
            "Services and Employees",
            "Academy (Courses, Students)",
            "Reports and AI Insights",
        ],
        "notes": [
            "Emphasize shared entities across modules.",
        ],
    },
    {
        "title": "POS Workflow",
        "bullets": [
            "Select or create customer",
            "Add services and products (scanner/manual)",
            "Apply discount and tax",
            "Choose payment mode",
            "Create bill and update stock/loyalty",
            "Print and optionally send WhatsApp receipt",
        ],
        "notes": [
            "POS speed and reliability are central to adoption.",
        ],
    },
    {
        "title": "Payment Flow",
        "bullets": [
            "Cash flow completes bill directly",
            "UPI/Card flow uses Razorpay order creation",
            "Callback signature is verified on backend",
            "Split payment fields are supported in bill schema",
        ],
        "notes": [
            "Security check is HMAC signature verification.",
        ],
    },
    {
        "title": "Appointments and Reminders",
        "bullets": [
            "Book appointment with customer, service, date, time",
            "Track status from scheduled to completion/cancellation",
            "Cron job runs daily at 8:00 AM",
            "WhatsApp reminders sent for next-day appointments",
        ],
        "notes": [
            "Point out reminderSent flag for duplicate prevention.",
        ],
    },
    {
        "title": "Inventory and Barcode",
        "bullets": [
            "Barcode lookup adds products quickly in POS",
            "Stock decremented on bill creation",
            "Low-stock endpoint helps reorder planning",
            "Stock restored when bill is cancelled",
        ],
        "notes": [
            "Data consistency between billing and inventory is a key strength.",
        ],
    },
    {
        "title": "Employee and Academy",
        "bullets": [
            "Employee setup includes role, salary, commission",
            "Performance endpoint aggregates billed contribution",
            "Academy supports courses, students, fee ledger, attendance",
            "Certificate issuance completes student lifecycle",
        ],
        "notes": [
            "App supports both retail services and education operations.",
        ],
    },
    {
        "title": "Reporting and KPIs",
        "bullets": [
            "Dashboard summary for daily and monthly health",
            "Sales trends with group-by periods",
            "Top services and products",
            "Payment method mix and employee performance",
        ],
        "notes": [
            "Reports drive pricing and staffing decisions.",
        ],
    },
    {
        "title": "AI Module",
        "bullets": [
            "Revenue prediction",
            "Churn-risk customer identification",
            "Service recommendations",
            "Peak-hour and inventory insights",
            "Rule-based chatbot endpoint for customer FAQs",
        ],
        "notes": [
            "Current AI is deterministic/statistical and can evolve later.",
        ],
    },
    {
        "title": "Security and Access Control",
        "bullets": [
            "JWT authentication for protected APIs",
            "Role control: admin vs staff",
            "Request validation via express-validator",
            "Helmet, CORS, and rate limiting globally enabled",
        ],
        "notes": [
            "Call out admin-only endpoints for critical operations.",
        ],
    },
    {
        "title": "Deployment Blueprint",
        "bullets": [
            "Frontend deployment on Vercel",
            "Backend deployment on Render or Railway",
            "Database on MongoDB Atlas",
            "Env-var driven secrets and integration setup",
        ],
        "notes": [
            "Mention post-deploy smoke test checklist.",
        ],
    },
    {
        "title": "Known Gaps and Roadmap",
        "bullets": [
            "Add appointment slot conflict prevention",
            "Align AI analytics fields with current schema",
            "Improve password policy and recovery flow",
            "Add centralized monitoring and audit logging",
        ],
        "notes": [
            "Present gaps as planned hardening, not system failure.",
        ],
    },
    {
        "title": "Closing and Next Steps",
        "bullets": [
            "Deliverables: workflow manual, full documentation, presentation deck",
            "Next: stakeholder review and UAT",
            "Then: staged rollout and production hardening",
        ],
        "notes": [
            "Invite operational and implementation questions.",
        ],
    },
]

prs = Presentation()

for item in slides:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = item["title"]

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for idx, bullet in enumerate(item["bullets"]):
        p = body.add_paragraph() if idx > 0 else body.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24 if idx == 0 else 20)

    notes_text = "\n".join(f"- {n}" for n in item["notes"])
    slide.notes_slide.notes_text_frame.text = notes_text

output_path = r"d:\Billing Software\galaxy-salon-system\docs\Galaxy-Salon-System-Presentation.pptx"
prs.save(output_path)
print(f"Presentation created: {output_path}")
